#!/usr/bin/env python3
"""
make_poster_pptx.py
===================
Generate an editable PowerPoint pipeline diagram for the QuadruJuggle poster.
Every box/arrow is a separate shape — drag freely in PowerPoint or LibreOffice Impress.

Output: docs/figures/poster_pipeline.pptx

Run:
    python scripts/make_poster_pptx.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import pptx.oxml.ns as nsmap
from lxml import etree

os.makedirs("docs/figures", exist_ok=True)

# ---------------------------------------------------------------------------
# Colour palette
# ---------------------------------------------------------------------------
C_BLUE   = RGBColor(0x19, 0x76, 0xD2)
C_ORANGE = RGBColor(0xF5, 0x7C, 0x00)
C_GREEN  = RGBColor(0x38, 0x8E, 0x3C)
C_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
C_GREY   = RGBColor(0x61, 0x61, 0x61)
C_RED    = RGBColor(0xC6, 0x28, 0x28)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_LTBLUE = RGBColor(0xE3, 0xF2, 0xFD)
C_LTGRN  = RGBColor(0xE8, 0xF5, 0xE9)
C_LTORG  = RGBColor(0xFF, 0xF3, 0xE0)
C_LTPUR  = RGBColor(0xF3, 0xE5, 0xF5)
C_LTGRY  = RGBColor(0xF5, 0xF5, 0xF5)

# ---------------------------------------------------------------------------
# Slide setup — widescreen 16:9 (13.33" × 7.5")
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.50)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)
shapes = slide.shapes


# ---------------------------------------------------------------------------
# Helper functions
# ---------------------------------------------------------------------------
def add_box(shapes, left, top, width, height,
            title, subtitle="",
            fill_color=C_LTBLUE, border_color=C_BLUE,
            title_color=C_BLUE, sub_color=C_GREY,
            title_size=Pt(12), sub_size=Pt(9),
            bold_title=True):
    """Add a rounded rectangle with title + optional subtitle text."""
    shape = shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.ROUNDED_RECTANGLE — use integer directly
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    # Override to rounded rectangle
    sp = shape._element
    spPr = sp.find(nsmap.qn("p:spPr"))
    prstGeom = spPr.find(nsmap.qn("a:prstGeom"))
    if prstGeom is not None:
        prstGeom.set("prst", "roundRect")
    else:
        prstGeom = etree.SubElement(spPr, nsmap.qn("a:prstGeom"))
        prstGeom.set("prst", "roundRect")
        avLst = etree.SubElement(prstGeom, nsmap.qn("a:avLst"))
        gd = etree.SubElement(avLst, nsmap.qn("a:gd"))
        gd.set("name", "adj"); gd.set("fmla", "val 30000")

    # Fill
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color

    # Border
    line = shape.line
    line.color.rgb = border_color
    line.width = Pt(1.8)

    # Text
    tf = shape.text_frame
    tf.word_wrap = True

    # Title paragraph
    p0 = tf.paragraphs[0]
    p0.alignment = PP_ALIGN.CENTER
    run0 = p0.add_run()
    run0.text = title
    run0.font.bold = bold_title
    run0.font.size = title_size
    run0.font.color.rgb = title_color

    # Subtitle paragraph
    if subtitle:
        from pptx.util import Pt as Pt2
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        run1 = p1.add_run()
        run1.text = subtitle
        run1.font.size = sub_size
        run1.font.color.rgb = sub_color
        run1.font.bold = False

    return shape


def add_label(shapes, left, top, width, height, text,
              color=C_GREY, size=Pt(9), bold=False, align=PP_ALIGN.CENTER):
    """Transparent text box."""
    txBox = shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox


def add_arrow(shapes, x1, y1, x2, y2, color=C_GREY, width=Pt(1.8)):
    """
    Add a straight connector (arrow) between two points.
    x1,y1 = start; x2,y2 = end, all in inches.
    """
    from pptx.util import Inches, Pt, Emu
    connector = shapes.add_connector(
        1,   # MSO_CONNECTOR_TYPE.STRAIGHT
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2),
    )
    connector.line.color.rgb = color
    connector.line.width = width
    # Add arrowhead
    ln = connector.line._ln
    tailEnd = ln.find(nsmap.qn("a:tailEnd"))
    if tailEnd is None:
        tailEnd = etree.SubElement(ln, nsmap.qn("a:tailEnd"))
    tailEnd.set("type", "none")
    headEnd = ln.find(nsmap.qn("a:headEnd"))
    if headEnd is None:
        headEnd = etree.SubElement(ln, nsmap.qn("a:headEnd"))
    headEnd.set("type", "arrow")
    headEnd.set("w", "med")
    headEnd.set("len", "med")
    return connector


# ---------------------------------------------------------------------------
# Title
# ---------------------------------------------------------------------------
add_label(shapes,
          left=0.3, top=0.10, width=12.5, height=0.55,
          text="QuadruJuggle — Hierarchical Control Architecture",
          color=C_BLACK, size=Pt(20), bold=True)

# ---------------------------------------------------------------------------
# Column positions (inches from left)
#   Col A: Sensing    x≈1.0
#   Col B: Pi1        x≈4.8  (wide box spanning two sub-cols)
#   Col C: Pi2        x≈8.9
#   Col D: Robot/Ball x≈11.4
# ---------------------------------------------------------------------------

# ── SENSING COLUMN ────────────────────────────────────────────────────────

add_box(shapes,
        left=0.25, top=1.0, width=2.2, height=1.1,
        title="Ball Sensor",
        subtitle="Camera / depth\nraw pos (noisy)",
        fill_color=C_LTGRY, border_color=C_GREY,
        title_color=C_GREY)

add_box(shapes,
        left=0.25, top=2.5, width=2.2, height=1.1,
        title="Kalman Filter",
        subtitle="position + velocity\nestimation @ 200 Hz",
        fill_color=C_LTPUR, border_color=C_PURPLE,
        title_color=C_PURPLE)

add_arrow(shapes, 1.35, 2.10, 1.35, 2.50, color=C_GREY)
add_label(shapes, 0.28, 2.12, 2.0, 0.30, "noisy pos", C_GREY, Pt(8))

# ── Pi1 OUTER GROUP BOX ───────────────────────────────────────────────────
# Dashed outer border — use a no-fill shape
outer = shapes.add_shape(1,
    Inches(2.80), Inches(1.0), Inches(5.5), Inches(3.0))
outer.fill.background()
outer.line.color.rgb = C_BLUE
outer.line.width = Pt(1.5)
outer.line.dash_style = 4   # DASH

# Pi1 header label
add_label(shapes,
          left=3.0, top=1.05, width=4.8, height=0.40,
          text="Pi1 — High-Level Planner  (choose one)",
          color=C_BLUE, size=Pt(11), bold=True)

# Mirror Law sub-box
add_box(shapes,
        left=2.95, top=1.55, width=2.4, height=1.65,
        title="Mirror Law",
        subtitle="Analytic geometry\nn̂ = norm(v_out/e − v_in)\nroll/pitch → 6D cmd\n(zero training needed)",
        fill_color=C_LTGRN, border_color=C_GREEN,
        title_color=C_GREEN, title_size=Pt(13))

# OR label
add_label(shapes,
          left=5.35, top=2.10, width=0.5, height=0.60,
          text="OR", color=C_GREY, size=Pt(13), bold=True)

# Learned Pi1 sub-box
add_box(shapes,
        left=5.90, top=1.55, width=2.2, height=1.65,
        title="Learned Pi1 (RL)",
        subtitle="PPO policy\n46D obs → 6D cmd\napex curriculum\n[0.10–0.60 m]",
        fill_color=C_LTORG, border_color=C_ORANGE,
        title_color=C_ORANGE, title_size=Pt(12))

# Switch logic label
add_box(shapes,
        left=3.20, top=3.30, width=4.8, height=0.58,
        title="Hybrid Switch: Pi1 builds energy  →  Mirror Law holds height",
        subtitle="",
        fill_color=RGBColor(0xE1, 0xF5, 0xFE), border_color=C_BLUE,
        title_color=C_BLUE, title_size=Pt(9.5), bold_title=False)

# Arrow from Kalman → Pi1 box
add_arrow(shapes, 2.45, 3.05, 2.82, 3.05, color=C_PURPLE)
add_label(shapes, 0.25, 3.08, 2.50, 0.38,
          "ball state (46D obs)", C_PURPLE, Pt(8.5))

# Arrows from mirror law / learned pi1 → combined output
add_arrow(shapes, 4.15, 3.20, 4.15, 3.88, color=C_GREEN)
add_arrow(shapes, 7.00, 3.20, 7.00, 3.60, color=C_ORANGE)
add_arrow(shapes, 7.00, 3.60, 4.45, 3.88, color=C_ORANGE)

# ── Pi2 BOX ───────────────────────────────────────────────────────────────
add_box(shapes,
        left=3.00, top=4.10, width=5.30, height=1.20,
        title="Pi2 — Torso Tracker  (frozen)",
        subtitle="Pre-trained RL policy  |  6D torso cmd → 12D joint targets  |  ~200M steps, never retrained",
        fill_color=C_LTBLUE, border_color=C_BLUE,
        title_color=C_BLUE, title_size=Pt(13))

add_arrow(shapes, 5.65, 3.90, 5.65, 4.10, color=C_BLUE)
add_label(shapes, 5.68, 3.90, 2.5, 0.30,
          "6D torso cmd  (h, ḣ, roll, pitch, vx, vy)", C_BLUE, Pt(8))

# ── ROBOT BOX ─────────────────────────────────────────────────────────────
add_box(shapes,
        left=9.10, top=1.80, width=3.80, height=1.30,
        title="Go1 Robot",
        subtitle="12-DOF quadruped\nPD control @ 200 Hz\ncustom paddle attachment",
        fill_color=C_LTGRY, border_color=C_GREY,
        title_color=C_GREY, title_size=Pt(13))

# Pi2 → Robot arrow (horizontal at row 4.7)
add_arrow(shapes, 8.30, 4.70, 11.50, 3.10, color=C_GREY)
add_label(shapes, 8.40, 3.70, 2.80, 0.30,
          "12D joint targets", C_GREY, Pt(8))

# ── BALL BOX ──────────────────────────────────────────────────────────────
add_box(shapes,
        left=9.60, top=3.50, width=2.80, height=1.00,
        title="Ball",
        subtitle="m = 50 g  |  e = 0.85\nphysical bouncing",
        fill_color=C_LTORG, border_color=C_ORANGE,
        title_color=C_ORANGE, title_size=Pt(13))

# Robot → Ball arrow
add_arrow(shapes, 11.00, 3.10, 11.00, 3.50, color=C_ORANGE)
add_label(shapes, 11.05, 3.15, 1.8, 0.30, "paddle impulse", C_ORANGE, Pt(8))

# Ball → Sensor feedback arrow (loop back left)
add_arrow(shapes, 9.60, 4.00, 1.45, 2.10, color=C_GREY)
add_label(shapes, 4.0, 4.80, 3.5, 0.30,
          "← ball trajectory (sensed)", C_GREY, Pt(8))

# ── TRAINING ANNOTATIONS ─────────────────────────────────────────────────
add_label(shapes, 2.95, 3.22, 2.45, 0.30,
          "✓ analytic — no training", C_GREEN, Pt(8))
add_label(shapes, 5.88, 3.22, 2.25, 0.30,
          "trained: PPO ~570M steps", C_ORANGE, Pt(8))
add_label(shapes, 3.00, 5.32, 5.30, 0.30,
          "trained: PPO ~200M steps — shared by all Pi1 variants", C_BLUE, Pt(8))

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out = "docs/figures/poster_pipeline.pptx"
prs.save(out)
print(f"Saved → {out}")
print("Open in PowerPoint or LibreOffice Impress. All boxes are independent shapes — drag freely.")
