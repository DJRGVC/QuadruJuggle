"""Generate the March 12 weekly capstone presentation.

Matches the exact styling of weekmarch3rdcapstone.pptx:
- Light theme (white/off-white backgrounds)
- Layout 11 (DEFAULT) for content slides, layout 2 (TITLE_AND_BODY) for title
- Arial 24pt bold #1F2937 for slide titles
- Calibri body text #6B7280 / #374151
- Pastel colored boxes with thin borders
- Horizontal connector lines #374151 2pt
- Legend squares at bottom
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Load template to inherit master/layouts/theme
prs = Presentation("slides/weekmarch3rdcapstone.pptx")

# Remove all existing slides
from lxml import etree
xml_slides = prs.slides._sldIdLst
for sldId in list(xml_slides):
    rId = sldId.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    if rId is None:
        rId = sldId.get('r:id')
    if rId:
        try:
            prs.part.drop_rel(rId)
        except KeyError:
            pass
    xml_slides.remove(sldId)

# -- Color palette (matches original) --
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)     # Titles, bold labels
GRAY_TEXT = RGBColor(0x6B, 0x72, 0x80)     # Subtitles, body, italic labels
MED_GRAY = RGBColor(0x37, 0x41, 0x51)      # Body text, lines
LIGHT_GRAY_BG = RGBColor(0xF3, 0xF4, 0xF6) # Gray box fill
VERY_LIGHT_BG = RGBColor(0xF9, 0xFA, 0xFB) # Content card fill
LINE_COLOR = RGBColor(0x37, 0x41, 0x51)     # Connector lines

# Box colors
YELLOW_FILL = RGBColor(0xFE, 0xF3, 0xC7)
YELLOW_BORDER = RGBColor(0xF5, 0x9E, 0x0B)
BLUE_FILL = RGBColor(0xDB, 0xEA, 0xFE)
BLUE_BORDER = RGBColor(0x3B, 0x82, 0xF6)
GREEN_FILL = RGBColor(0xD1, 0xFA, 0xE5)
GREEN_BORDER = RGBColor(0x10, 0xB9, 0x81)
GRAY_BORDER = RGBColor(0x9C, 0xA3, 0xAF)
PURPLE_FILL = RGBColor(0xED, 0xEA, 0xFE)
PURPLE_BORDER = RGBColor(0x63, 0x66, 0xF1)
RED_FILL = RGBColor(0xFE, 0xE2, 0xE2)
RED_BORDER = RGBColor(0xEF, 0x44, 0x44)
ORANGE_FILL = RGBColor(0xFF, 0xED, 0xD5)
ORANGE_BORDER = RGBColor(0xF9, 0x73, 0x16)
DARK_BORDER = RGBColor(0x33, 0x33, 0x33)

SW = prs.slide_width   # 9144000 EMU
SH = prs.slide_height  # 5143500 EMU

# -- Font sizes in EMU (matching original exactly) --
TITLE_SIZE = 304800     # 24pt
SUBTITLE_SIZE = 165100  # 13pt
BODY_SIZE = 139700      # 11pt
LABEL_SIZE = 127000     # 10pt
SMALL_SIZE = 120650     # 9.5pt
CAPTION_SIZE = 114300   # 9pt
BOX_TITLE_SIZE = 203200 # 16pt
PI_SIZE = 254000        # 20pt
MEDIUM_SIZE = 152400    # 12pt


def add_slide_default():
    """Add a slide using DEFAULT layout (idx 11) — matches original content slides."""
    layout = prs.slide_layouts[11]
    return prs.slides.add_slide(layout)


def add_slide_title():
    """Add a slide using TITLE_AND_BODY layout (idx 2) — matches original title slide."""
    layout = prs.slide_layouts[2]
    return prs.slides.add_slide(layout)


def add_title(slide, text, left=457200, top=228600, width=8229600, height=502800):
    """Add slide title matching original: Arial 24pt bold #1F2937."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                    Emu(width), Emu(height))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Arial"
    run.font.size = Emu(TITLE_SIZE)
    run.font.bold = True
    run.font.color.rgb = DARK_TEXT
    return shape


def add_label(slide, text, left, top, width=1188600, height=228600,
              font_size=LABEL_SIZE, color=GRAY_TEXT, bold=False, italic=True,
              alignment=PP_ALIGN.CENTER, font_name="Arial"):
    """Add an italic gray label (column headers, annotations)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                    Emu(width), Emu(height))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Emu(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    return shape


def add_box(slide, left, top, width, height, fill_color, border_color, border_width=19050):
    """Add a colored rectangle with border (matching original box style)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                    Emu(width), Emu(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Emu(border_width)
    return shape


def add_box_text(slide, text, left, top, width, height,
                 font_size=BOX_TITLE_SIZE, color=DARK_TEXT, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name="Arial"):
    """Add text overlay on a box (transparent bg)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                    Emu(width), Emu(height))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Emu(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = alignment
    return shape


def add_line(slide, left, top, width, height=0):
    """Add a connector line matching original: #374151, 2pt."""
    connector = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                        Emu(left), Emu(top),
                                        Emu(max(width, 12700)),
                                        Emu(max(height, 12700)))
    connector.fill.solid()
    connector.fill.fore_color.rgb = LINE_COLOR
    connector.line.fill.background()
    return connector


def add_body_text(slide, left, top, width, height, paragraphs, font_size=BODY_SIZE):
    """Add multi-paragraph body text. paragraphs is list of (text, color, bold, italic)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(left), Emu(top),
                                    Emu(width), Emu(height))
    shape.fill.background()
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(paragraphs):
        if isinstance(item, str):
            text, color, bold, italic = item, GRAY_TEXT, False, False
        elif len(item) == 2:
            text, color = item
            bold, italic = False, False
        elif len(item) == 3:
            text, color, bold = item
            bold, italic = bold, False
        else:
            text, color, bold, italic = item
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = text
        run.font.name = "Calibri"
        run.font.size = Emu(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(1)
        p.line_spacing = Pt(font_size / 914400 * 72 + 2)  # font size in pt + 2pt leading
    return shape


def add_content_card(slide, left, top, width, height, title, title_color, body_lines,
                     bg_color=VERY_LIGHT_BG, border_color=DARK_BORDER):
    """Add a content card (like the prior work cards in original slide 5)."""
    add_box(slide, left, top, width, height, bg_color, border_color, 12700)
    inner_left = left + 91440
    inner_top = top + 45720
    inner_w = width - 182880
    # Title
    add_box_text(slide, title, inner_left, inner_top, inner_w, 182880,
                 font_size=MEDIUM_SIZE, color=title_color, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    # Body lines
    paras = [(line, MED_GRAY, False, False) for line in body_lines]
    add_body_text(slide, inner_left, inner_top + 182880, inner_w, height - 228600,
                  paras, font_size=SMALL_SIZE)


def add_legend(slide, top, items):
    """Add a legend row. items is list of (fill_color, border_color, label)."""
    x = 457200
    for fill, border, label in items:
        add_box(slide, x, top, 228600, 228600, fill, border, 12700)
        add_label(slide, label, x + 274320, top, 914400, 228600,
                  font_size=CAPTION_SIZE, color=GRAY_TEXT, bold=False, italic=False,
                  alignment=PP_ALIGN.LEFT)
        x += 1371600


# ============================================================
# SLIDE 1: Title (TITLE_AND_BODY layout, like original)
# ============================================================
s = add_slide_title()
# The TITLE_AND_BODY layout has a placeholder — fill it
for shape in s.placeholders:
    if shape.placeholder_format.idx == 0:  # title placeholder
        shape.text = "MEng 164: Physics-Informed Continual Learning for Stochastic Optimal Control of Legged Robots"
        break

# Subtitle
add_label(s, "Graduate Mentor: Sangli Teng, Faculty Advisor: Prof. Koushil Sreenath",
          490500, 1404950, 8028000, 461700, font_size=228600, color=GRAY_TEXT,
          italic=False, alignment=PP_ALIGN.LEFT)

# Names
add_label(s, "Frank Wang", 1250400, 4514825, 1605300, 461700,
          font_size=228600, color=GRAY_TEXT, italic=False, alignment=PP_ALIGN.LEFT)
add_label(s, "Daniel Grant", 3769350, 4514825, 1605300, 461700,
          font_size=228600, color=GRAY_TEXT, italic=False, alignment=PP_ALIGN.LEFT)
add_label(s, "Jaime de Carlos", 6288300, 4514825, 1995300, 461700,
          font_size=228600, color=GRAY_TEXT, italic=False, alignment=PP_ALIGN.LEFT)

# Week indicator
add_label(s, "Weekly Update — March 12, 2026", 490500, 2000000, 8028000, 461700,
          font_size=203200, color=BLUE_BORDER, italic=False, bold=True,
          alignment=PP_ALIGN.LEFT, font_name="Calibri")

add_label(s, "Pi2 Walking Breakthrough  +  EKF Perception Roadmap",
          490500, 2700000, 8028000, 461700,
          font_size=177800, color=MED_GRAY, italic=True, bold=False,
          alignment=PP_ALIGN.LEFT, font_name="Calibri")


# ============================================================
# SLIDE 2: Updated Architecture Diagram (matches slide 7 of original)
# ============================================================
s = add_slide_default()
add_title(s, "Updated Hierarchical Architecture")

# Column labels (italic gray, matching original)
add_label(s, "Input", 457200, 914400)
add_label(s, "Ball Planner", 2011680, 914400, 1554600)
add_label(s, "8D Interface", 3977640, 914400, 1828800)
add_label(s, "Torso Tracker", 6035040, 914400, 1554600)
add_label(s, "Robot", 7955280, 914400, 914400)

# Ball input box (yellow, like original)
add_box(s, 457200, 1234440, 1188600, 822900, YELLOW_FILL, YELLOW_BORDER)
add_box_text(s, "Ball", 457200, 1234440, 1188600, 500000)
add_label(s, "p, ṗ ∈ ℝ⁶", 457200, 2103120, 1188600, 228600)
add_label(s, "(EKF in future)", 457200, 2331720, 1188600, 228600)

# Arrow: Ball → π₁
add_line(s, 1645920, 1645920, 365700)

# π₁ box (blue, like original)
add_box(s, 2011680, 1234440, 1554600, 822900, BLUE_FILL, BLUE_BORDER)
add_box_text(s, "π₁", 2011680, 1234440, 1554600, 822900, font_size=PI_SIZE)

# Arrow: π₁ → 8D
add_line(s, 3566280, 1645920, 411360)

# 8D interface box (orange/amber)
add_box(s, 3977640, 1234440, 1828800, 822900, ORANGE_FILL, ORANGE_BORDER)
add_box_text(s, "8D cmd", 3977640, 1280000, 1828800, 300000, font_size=BOX_TITLE_SIZE)
# 8D details inside box
add_body_text(s, 4050000, 1500000, 1700000, 550000, [
    ("h  ḣ  roll  pitch", MED_GRAY, False, False),
    ("ωᵣ  ωₚ  vₓ  vᵧ", MED_GRAY, False, False),
], font_size=CAPTION_SIZE)

# Arrow: 8D → π₂
add_line(s, 5806440, 1645920, 228600)

# π₂ box (green, like original robot box)
add_box(s, 6035040, 1234440, 1554600, 822900, GREEN_FILL, GREEN_BORDER)
add_box_text(s, "π₂", 6035040, 1234440, 1554600, 822900, font_size=PI_SIZE)

# Arrow: π₂ → Robot
add_line(s, 7589640, 1645920, 365640)

# Robot box (green)
add_box(s, 7955280, 1234440, 914400, 822900, GREEN_FILL, GREEN_BORDER)
add_box_text(s, "Go1", 7955280, 1234440, 914400, 822900, font_size=BOX_TITLE_SIZE)

# Feedback line: Robot → π₂ → π₁ (horizontal at bottom)
add_line(s, 3200400, 2057340, 12700, 685800)  # vertical from π₁
add_line(s, 3200400, 2743200, 5206680)         # horizontal
add_line(s, 8407080, 2057400, 12700, 685800)  # vertical from Robot

# Robot State label
add_label(s, "Robot State", 5029200, 2788920, 2743200, 228600)

# qᵈₘ label under π₂ → Robot
add_label(s, "qᵈₘ ∈ ℝ¹²", 6035040, 2103120, 1554600, 228600)
add_label(s, "12 joint targets", 6035040, 2300000, 1554600, 228600, font_size=CAPTION_SIZE)

# Legend
add_legend(s, 3100000, [
    (BLUE_FILL, BLUE_BORDER, "RL Policy"),
    (GREEN_FILL, GREEN_BORDER, "Robot"),
    (YELLOW_FILL, YELLOW_BORDER, "Input / State"),
    (ORANGE_FILL, ORANGE_BORDER, "Interface"),
])

# Bottom description
add_body_text(s, 457200, 3500000, 8229600, 914400, [
    ("π₁ (Ball Planner): ", DARK_TEXT, True, False),
    ("  Ball state + proprio → 8D torso commands. MLP [256, 128, 64]. Trained 2nd (frozen π₂).", GRAY_TEXT),
    ("π₂ (Torso Tracker): ", DARK_TEXT, True, False),
    ("  8D cmd + proprio (53D obs) → 12 joint targets. MLP [512, 256, 128]. Trained 1st, standalone.", GRAY_TEXT),
    ("8D Interface: ", DARK_TEXT, True, False),
    ("  [h, ḣ, roll, pitch, ωᵣ, ωₚ, vₓ, vᵧ]. Height [0.20, 0.50]m, vxy [−0.50, 0.50]m/s, ω [−4.0, 4.0]rad/s", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 3: Pi2 — What Went Wrong
# ============================================================
s = add_slide_default()
add_title(s, "Pi2 Training: What Went Wrong")

# Left card: The problem
add_content_card(s, 457200, 914400, 3931800, 3840600,
                 "The Problem: Robot Was a Statue",
                 RED_BORDER,
                 [
                     "Standing still earned higher reward than walking.",
                     "",
                     "foot_contact penalty: −0.5 per foot airborne.",
                     "  → Walking lifts 2 feet = −1.0/step > velocity reward.",
                     "",
                     "vxy_error penalty defined in code but never wired",
                     "  to RewardsCfg! (reward function existed, weight = 0)",
                     "",
                     "feet_air_time = 0.01 (Isaac Lab uses 0.25)",
                     "",
                     "PPO exploration too low:",
                     "  init_noise_std = 0.5 (Isaac Lab: 1.0)",
                     "  entropy_coef = 0.001 (Isaac Lab: 0.01)",
                     "  Network [256, 128] (Isaac Lab: [512, 256, 128])",
                     "",
                     "Missing last_action observation (12D)",
                     "Aggressive DR prevented gait discovery",
                 ],
                 bg_color=RED_FILL, border_color=RED_BORDER)

# Right card: The fix
add_content_card(s, 4754880, 914400, 3931800, 3840600,
                 "The Fix: Align with Isaac Lab Go1 Defaults",
                 GREEN_BORDER,
                 [
                     "Removed foot_contact penalty entirely",
                     "Wired vxy_error penalty at weight −4.0",
                     "feet_air_time: 0.01 → 0.25",
                     "",
                     "init_noise_std: 0.5 → 1.0",
                     "entropy_coef: 0.001 → 0.01",
                     "Network: [256, 128] → [512, 256, 128]",
                     "num_steps_per_env: 24 → 48",
                     "",
                     "Added last_action to observations (53D total)",
                     "Stripped DR to Isaac Lab minimum",
                     "Command resampling: 2−5s → 5−10s",
                     "",
                     "Key insight: not just reward shaping — it was",
                     "anti-locomotion penalties + insufficient exploration",
                     "+ missing observations + aggressive DR, all",
                     "compounding to make walking impossible.",
                 ],
                 bg_color=GREEN_FILL, border_color=GREEN_BORDER)


# ============================================================
# SLIDE 4: Pi2 — 3-Stage Curriculum
# ============================================================
s = add_slide_default()
add_title(s, "Pi2: 3-Stage Walk-First Curriculum")

add_body_text(s, 457200, 800000, 8229600, 300000, [
    ("Simplified from 8 stages to 3. Walk first, then add pose tracking.", GRAY_TEXT),
], font_size=BODY_SIZE)

# Stage A box
add_box(s, 457200, 1150000, 8229600, 700000, BLUE_FILL, BLUE_BORDER, 12700)
add_body_text(s, 570000, 1165000, 7900000, 670000, [
    ("Stage A — Walk First", DARK_TEXT, True, False),
    ("Fixed pose (h = 0.38m, no tilt, no rotation). Full velocity range vxy = [−0.50, 0.50] m/s.", MED_GRAY),
    ("Advancement gate: timeout ≥ 85%  AND  vxy_error > −0.80  (proves actual walking, not standing still)", MED_GRAY),
], font_size=SMALL_SIZE)

# Stage B box
add_box(s, 457200, 1930000, 8229600, 700000, PURPLE_FILL, PURPLE_BORDER, 12700)
add_body_text(s, 570000, 1945000, 7900000, 670000, [
    ("Stage B — Mild Pose + Walking", DARK_TEXT, True, False),
    ("h = [0.32, 0.44]m, roll/pitch = ±0.15 rad, ω = ±1.0 rad/s. Velocity kept at ±0.50 m/s.", MED_GRAY),
    ("Advancement gate: timeout ≥ 85% for 100 consecutive iterations", MED_GRAY),
], font_size=SMALL_SIZE)

# Stage C box
add_box(s, 457200, 2710000, 8229600, 700000, ORANGE_FILL, ORANGE_BORDER, 12700)
add_body_text(s, 570000, 2725000, 7900000, 670000, [
    ("Stage C — Full 8D", DARK_TEXT, True, False),
    ("h = [0.20, 0.50]m, roll/pitch = ±0.50 rad, ω = ±4.0 rad/s, vxy = ±0.50 m/s.", MED_GRAY),
    ("Early stopping: patience = 500 iterations. Currently training (iter 294/5000).", MED_GRAY),
], font_size=SMALL_SIZE)

# Bottom note
add_body_text(s, 457200, 3550000, 8229600, 300000, [
    ("Training config: ", DARK_TEXT, True, False),
    ("12,288 envs | PPO | [512, 256, 128] MLP | Isaac Lab + RSL-RL | ~5.4s/iter | ETA ~7 hours", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 5: Pi2 Training Results
# ============================================================
s = add_slide_default()
add_title(s, "Pi2 Results: Stage C (Full 8D), Iteration 294")

# Left: Metrics card
add_content_card(s, 457200, 914400, 3931800, 3840600,
                 "Training Metrics",
                 BLUE_BORDER,
                 [
                     "vx_tracking reward:        1.14",
                     "vy_tracking reward:        1.15",
                     "vxy_error penalty:          −0.32",
                     "",
                     "height_error:                −0.057",
                     "roll_error:                    −0.029",
                     "pitch_error:                  −0.031",
                     "roll_rate_error:             −0.258",
                     "pitch_rate_error:           −0.161",
                     "",
                     "Timeout:                        99.1%",
                     "Trunk collapsed:             0.9%",
                     "Mean reward:                 13.69",
                     "Noise std:                      0.93",
                     "Early stop counter:         6 / 500",
                 ],
                 bg_color=BLUE_FILL, border_color=BLUE_BORDER)

# Right: Interpretation card
add_content_card(s, 4754880, 914400, 3931800, 3840600,
                 "Interpretation",
                 GREEN_BORDER,
                 [
                     "Robot IS walking — vx/vy tracking both positive.",
                     "",
                     "99% timeout = survives full 20s episodes.",
                     "Only 0.9% trunk collapse = stable gait.",
                     "",
                     "Reward dropped from ~24 (Stage B) to ~14 (Stage C).",
                     "  Expected: full 8D is harder than walking alone.",
                     "",
                     "ES counter 6/500 = reward still improving.",
                     "",
                     "Rate errors (ωᵣ, ωₚ) are largest — hardest to track.",
                     "Position errors (h, roll, pitch) are small.",
                     "",
                     "Training progressed A → B → C without manual",
                     "  intervention. Stage A gate (vxy > −0.80)",
                     "  prevented fake advancement by standing still.",
                 ],
                 bg_color=GREEN_FILL, border_color=GREEN_BORDER)


# ============================================================
# SLIDE 6: Perception Architecture Change
# ============================================================
s = add_slide_default()
add_title(s, "Perception: EKF, Not Teacher-Student")

add_body_text(s, 457200, 731520, 8229600, 228600, [
    ("Last week: CNN encoder + teacher-student distillation (slide 11). New approach below.", GRAY_TEXT, False, True),
], font_size=BODY_SIZE)

# Left card: Why not CNN
add_content_card(s, 457200, 1100000, 3931800, 2400000,
                 "Why NOT Teacher-Student Distillation",
                 RED_BORDER,
                 [
                     "Rendering cameras for 12,288 envs = 10−50× slower",
                     "Visual domain gap (sim images ≠ real images)",
                     "Need visual DR (lighting, shadows, backgrounds)",
                     "Single frame can't give velocity — need temporal CNN",
                     "Creates a NEW sim-to-real gap to close",
                     "CNN architecture design + training pipeline overhead",
                 ],
                 bg_color=RED_FILL, border_color=RED_BORDER)

# Right card: Why EKF
add_content_card(s, 4754880, 1100000, 3931800, 2400000,
                 "Why EKF (Ma et al., Sci. Robotics 2025)",
                 GREEN_BORDER,
                 [
                     "Inject 6-number noise model instead of rendering",
                     "Deterministic code — identical in sim and real",
                     "Zero visual domain gap",
                     "Physics are known: ballistic + drag = 1 equation",
                     "~50 FLOPs/step, <1ms on Jetson Nano",
                     "ETH result: ~100% clean, 60−80% with real noise",
                 ],
                 bg_color=GREEN_FILL, border_color=GREEN_BORDER)

# Pipeline diagram boxes
y_pipe = 3700000
add_label(s, "Sim training:", 457200, y_pipe, 1371600, 228600,
          font_size=SMALL_SIZE, color=DARK_TEXT, bold=True, italic=False,
          alignment=PP_ALIGN.LEFT, font_name="Calibri")
# GT box
add_box(s, 1828800, y_pipe, 1097400, 228600, YELLOW_FILL, YELLOW_BORDER, 12700)
add_box_text(s, "GT + noise", 1828800, y_pipe, 1097400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT)
add_line(s, 2926200, y_pipe + 100000, 182880)
# EKF box
add_box(s, 3109080, y_pipe, 914400, 228600, PURPLE_FILL, PURPLE_BORDER, 12700)
add_box_text(s, "EKF", 3109080, y_pipe, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT, bold=True)
add_line(s, 4023480, y_pipe + 100000, 182880)
# π₁ box
add_box(s, 4206360, y_pipe, 914400, 228600, BLUE_FILL, BLUE_BORDER, 12700)
add_box_text(s, "π₁ → 8D", 4206360, y_pipe, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT)
add_line(s, 5120760, y_pipe + 100000, 182880)
# π₂ box
add_box(s, 5303640, y_pipe, 914400, 228600, GREEN_FILL, GREEN_BORDER, 12700)
add_box_text(s, "π₂ → 12D", 5303640, y_pipe, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT)
add_line(s, 6218040, y_pipe + 100000, 182880)
# Go1 box
add_box(s, 6400920, y_pipe, 914400, 228600, GREEN_FILL, GREEN_BORDER, 12700)
add_box_text(s, "Go1", 6400920, y_pipe, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT, bold=True)

# Real deploy row
y_real = 4100000
add_label(s, "Real deploy:", 457200, y_real, 1371600, 228600,
          font_size=SMALL_SIZE, color=DARK_TEXT, bold=True, italic=False,
          alignment=PP_ALIGN.LEFT, font_name="Calibri")
add_box(s, 1828800, y_real, 1097400, 228600, YELLOW_FILL, YELLOW_BORDER, 12700)
add_box_text(s, "Camera + HSV", 1828800, y_real, 1097400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT)
add_line(s, 2926200, y_real + 100000, 182880)
add_box(s, 3109080, y_real, 914400, 228600, PURPLE_FILL, PURPLE_BORDER, 12700)
add_box_text(s, "EKF", 3109080, y_real, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT, bold=True)
add_line(s, 4023480, y_real + 100000, 182880)
add_box(s, 4206360, y_real, 914400, 228600, BLUE_FILL, BLUE_BORDER, 12700)
add_box_text(s, "π₁ → 8D", 4206360, y_real, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT)
add_line(s, 5120760, y_real + 100000, 182880)
add_box(s, 5303640, y_real, 914400, 228600, GREEN_FILL, GREEN_BORDER, 12700)
add_box_text(s, "π₂ → 12D", 5303640, y_real, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT)
add_line(s, 6218040, y_real + 100000, 182880)
add_box(s, 6400920, y_real, 914400, 228600, GREEN_FILL, GREEN_BORDER, 12700)
add_box_text(s, "Go1", 6400920, y_real, 914400, 228600,
             font_size=CAPTION_SIZE, color=DARK_TEXT, bold=True)

# "identical code" annotation pointing to EKF
add_label(s, "← identical code in sim & real", 4023480, y_real + 280000, 3000000, 200000,
          font_size=CAPTION_SIZE, color=BLUE_BORDER, bold=True, italic=True,
          alignment=PP_ALIGN.LEFT, font_name="Calibri")

# Key advantage
add_body_text(s, 457200, 4650000, 8229600, 300000, [
    ("Hierarchy advantage: ", DARK_TEXT, True, False),
    ("only π₁ needs perception-aware retraining. π₂ never sees ball state — completely isolated from perception.", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 7: EKF Specification
# ============================================================
s = add_slide_default()
add_title(s, "EKF Specification")

# Left card: Predict
add_content_card(s, 457200, 914400, 3931800, 1600000,
                 "Predict Step (200 Hz, every physics step)",
                 BLUE_BORDER,
                 [
                     "State: x̂ = [pₓ, pᵧ, p_z, vₓ, vᵧ, v_z]  (6D, paddle frame)",
                     "",
                     "Ballistic dynamics + quadratic air drag:",
                     "  drag = 0.114 s⁻¹m⁻¹ (40mm ping-pong ball)",
                     "  a = [0, 0, −9.81] − drag · |v| · v",
                     "  v_new = v + a · dt    (dt = 0.005s)",
                     "  p_new = p + v · dt + 0.5 · a · dt²",
                 ],
                 bg_color=BLUE_FILL, border_color=BLUE_BORDER)

# Left card: Update
add_content_card(s, 457200, 2700000, 3931800, 1400000,
                 "Update Step (60−120 Hz, when camera sees ball)",
                 PURPLE_BORDER,
                 [
                     "z = [pₓ, pᵧ, p_z] from detector",
                     "Innovation: y = z − H · x̂   (H = [I₃ | 0₃])",
                     "Kalman gain: K = P · Hᵀ · S⁻¹",
                     "",
                     "Missed frame? Skip update, coast on ballistic model.",
                     "  Accurate to ~2mm over 50−100ms for a ping-pong ball.",
                 ],
                 bg_color=PURPLE_FILL, border_color=PURPLE_BORDER)

# Right card: Detection pipeline
add_content_card(s, 4754880, 914400, 3931800, 2200000,
                 "Detection Pipeline (real robot)",
                 ORANGE_BORDER,
                 [
                     "1. Camera frame (120 fps, global shutter USB)",
                     "2. BGR → HSV threshold (orange ball)         ~0.3ms",
                     "3. Contours → largest → centroid (u, v) + diam   ~0.2ms",
                     "4. Monocular depth: z = 0.040 · f / d_px        ~0 ms",
                     "5. Back-project to 3D:  x = (u−cₓ)·z/fₓ         ~0 ms",
                     "6. Transform to paddle frame (fixed rigid)     ~0 ms",
                     "",
                     "Total compute: < 1 ms on Jetson",
                     "No neural network. No learned parameters.",
                     "Just geometry and one division.",
                 ],
                 bg_color=ORANGE_FILL, border_color=ORANGE_BORDER)

# Noise model card
add_content_card(s, 4754880, 3300000, 3931800, 900000,
                 "Noise Model (placeholder params)",
                 GRAY_BORDER,
                 [
                     "σ = σ_base(3mm) + σ_dist(5mm/m)·d + σ_ω(2mm/(rad/s))·|ω|",
                     "Detection dropout: 5−15% per frame",
                     "Asymmetric AC: actor = EKF output, critic = ground truth",
                 ],
                 bg_color=LIGHT_GRAY_BG, border_color=GRAY_BORDER)

# Implementation note
add_body_text(s, 457200, 4350000, 8229600, 400000, [
    ("Implementation: ", DARK_TEXT, True, False),
    ("pure PyTorch, GPU-batched for 12,288 envs. Zero learned params. ~50 FLOPs/step.", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 8: 6-Phase Roadmap (table-like, matching original's structure)
# ============================================================
s = add_slide_default()
add_title(s, "6-Phase Perception Pipeline")

phases = [
    ("Phase 0", "Privileged Training (current)", "IN PROGRESS",
     GREEN_FILL, GREEN_BORDER, "Pi2 walking ✓. Pi1 training next."),
    ("Phase 1", "EKF in Sim (clean data)", "Week 2",
     BLUE_FILL, BLUE_BORDER, "Validate filter on pi1's ball trajectories"),
    ("Phase 2", "Noise Model + Retrain π₁", "Week 3",
     BLUE_FILL, BLUE_BORDER, "Perception-aware training, asymmetric actor-critic"),
    ("Phase 3", "Hardware Setup", "Weeks 1−2 (||)",
     ORANGE_FILL, ORANGE_BORDER, "Camera, bracket, paddle, cable routing on Go1"),
    ("Phase 4", "Real Detection Pipeline", "Weeks 2−3",
     ORANGE_FILL, ORANGE_BORDER, "HSV + monocular depth running on Jetson"),
    ("Phase 5", "Noise Calibration", "Weeks 3−4",
     ORANGE_FILL, ORANGE_BORDER, "Real data → regress σ(d, ω). ≥500 measurements"),
    ("Phase 6", "Real Deployment", "Weeks 5−6",
     RED_FILL, RED_BORDER, "Progressive: comms → π₂ → π₂+π₁ → perception → juggling"),
]

y = 900000
row_h = 500000
for phase_label, name, when, fill, border, desc in phases:
    add_box(s, 457200, y, 1188600, 380000, fill, border, 12700)
    add_box_text(s, phase_label, 457200, y, 1188600, 380000,
                 font_size=LABEL_SIZE, color=DARK_TEXT, bold=True)
    add_box_text(s, name, 1828800, y + 20000, 2743200, 200000,
                 font_size=BODY_SIZE, color=DARK_TEXT, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    add_label(s, when, 1828800, y + 200000, 2743200, 180000,
              font_size=CAPTION_SIZE, color=GRAY_TEXT, bold=False, italic=True,
              alignment=PP_ALIGN.LEFT, font_name="Calibri")
    add_label(s, desc, 4754880, y + 50000, 4100000, 300000,
              font_size=SMALL_SIZE, color=MED_GRAY, bold=False, italic=False,
              alignment=PP_ALIGN.LEFT, font_name="Calibri")
    y += row_h

# Note at bottom
add_body_text(s, 457200, y + 100000, 8229600, 400000, [
    ("Phases 1−2 (sim) and 3−5 (hardware) run on parallel tracks. Phase 6 merges both.", MED_GRAY, False, True),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 9: Sim-to-Real Gaps Diagram
# ============================================================
s = add_slide_default()
add_title(s, "Five Sim-to-Real Gaps")

gaps = [
    ("Actuator Dynamics", "Sim motors ideal; real have friction, delay, bandwidth",
     "System ID (CMA-ES) + domain randomization",
     ORANGE_FILL, ORANGE_BORDER),
    ("Perception", "Sim has ground-truth; real has noisy camera",
     "EKF + noise model (ETH approach, Ma et al. 2025)",
     PURPLE_FILL, PURPLE_BORDER),
    ("Latency", "Sim is synchronous; real has ~15ms pipeline delay",
     "Obs delay (1 policy step) + action delay (1 physics step)",
     BLUE_FILL, BLUE_BORDER),
    ("Physics", "Contact, friction, restitution differ from real",
     "Domain randomization + ball physics calibration",
     YELLOW_FILL, YELLOW_BORDER),
    ("Mechanical", "Sim robot is rigid; real has flex, backlash",
     "System ID + mechanical robustness testing",
     LIGHT_GRAY_BG, GRAY_BORDER),
]

y = 900000
for gap_name, what, how, fill, border in gaps:
    add_box(s, 457200, y, 8229600, 640000, fill, border, 12700)
    add_box_text(s, gap_name, 640000, y + 30000, 2100000, 250000,
                 font_size=MEDIUM_SIZE, color=DARK_TEXT, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    add_label(s, what, 640000, y + 300000, 3500000, 250000,
              font_size=SMALL_SIZE, color=MED_GRAY, italic=False,
              alignment=PP_ALIGN.LEFT, font_name="Calibri")
    add_label(s, how, 4300000, y + 130000, 4200000, 300000,
              font_size=SMALL_SIZE, color=DARK_TEXT, italic=False, bold=False,
              alignment=PP_ALIGN.LEFT, font_name="Calibri")
    y += 720000

add_body_text(s, 457200, y + 100000, 8229600, 500000, [
    ("Highest risk: ", DARK_TEXT, True, False),
    ("Go1 actuator bandwidth. Geared motors may not track fast corrections needed for juggling. "
     "Mitigation: sine sweep characterization early. Potentially reduce policy freq 50 → 25 Hz.", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 10: Key References (matches original slide 5 layout)
# ============================================================
s = add_slide_default()
add_title(s, "Key References")

# Left card
add_content_card(s, 457200, 914400, 3931800, 2286000,
                 "Ma et al. — ETH Zurich, Sci. Robotics 2025",
                 BLUE_BORDER,
                 [
                     "\"Learning coordinated badminton skills for legged",
                     "manipulators.\"",
                     "",
                     "Primary approach for our perception pipeline:",
                     "• EKF with noise-injected training",
                     "• Asymmetric actor-critic (noisy actor, clean critic)",
                     "• Result: ~100% clean, 60−80% with real noise",
                     "• Near-zero-shot sim-to-real on ANYmal + DynaArm",
                 ],
                 bg_color=VERY_LIGHT_BG, border_color=DARK_BORDER)

# Right card
add_content_card(s, 4754880, 914400, 3931800, 2286000,
                 "Su et al. — Berkeley HITTER, 2025",
                 BLUE_BORDER,
                 [
                     "Hierarchical planner + controller in Isaac Lab.",
                     "",
                     "Validates our hierarchical architecture:",
                     "• Pi1 (planner) → Pi2 (controller) decomposition",
                     "• PPO training with curriculum in Isaac Lab",
                     "• Zero-shot sim-to-real on Unitree G1",
                     "• Similar decoupled training: freeze controller,",
                     "  train planner on top",
                 ],
                 bg_color=VERY_LIGHT_BG, border_color=DARK_BORDER)

# Bottom references
add_body_text(s, 457200, 3400000, 8229600, 1500000, [
    ("Additional references reviewed:", DARK_TEXT, True, False),
    ("D'Ambrosio et al. (DeepMind, 2024) — Latency modeling \"crucial\"; 27k-param CNN on raw Bayer", MED_GRAY),
    ("Huang et al. (IROS 2023, KeeperBot) — Hierarchical RL for quadruped goalkeeping (UC Berkeley)", MED_GRAY),
    ("Ji, Margolis, Agrawal (ICRA 2023, DribbleBot) — End-to-end RL, DR, foot contact rewards (MIT)", MED_GRAY),
    ("Rudin et al. (CoRL 2022, legged_gym) — Standard DR + RSL-RL for quadrupeds (ETH)", MED_GRAY),
    ("Margolis & Agrawal (RSS 2022, Walk-These-Ways) — Motor gain DR, velocity tracking (MIT)", MED_GRAY),
    ("Ziegler et al. (Tübingen, 2025) — Event cameras 27.8× faster than frame-based (future option)", MED_GRAY),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 11: Next Steps & Timeline (matches original slide 12)
# ============================================================
s = add_slide_default()
add_title(s, "Next Steps & Timeline")

# Phase boxes (matching original slide 12 style)
phases_next = [
    ("This Week", "Finish Pi2 Training + Record Video", GREEN_FILL, GREEN_BORDER,
     "Stage C training running (~7h ETA). Record walking + velocity tracking demo. "
     "Then: re-enable DR (mass, friction, pushes) and retrain pi2."),
    ("Next Week", "Pi1 Training + EKF Build", BLUE_FILL, BLUE_BORDER,
     "Train pi1 (ball planner) on top of DR-hardened pi2. "
     "Build + validate batched EKF in PyTorch on sim trajectories."),
    ("Week 3", "Perception-Aware Pi1 + Hardware Setup", PURPLE_FILL, PURPLE_BORDER,
     "Retrain pi1 with noisy EKF observations (asymmetric AC). "
     "Order camera, 3D-print bracket, mount paddle on Go1."),
    ("Weeks 4−6", "System ID + Progressive Deployment", ORANGE_FILL, ORANGE_BORDER,
     "Actuator sine sweeps, ball bounce calibration. "
     "Deploy: comms check → pi2 walking → ball balancing → juggling."),
]

y = 914400
for label, name, fill, border, desc in phases_next:
    add_box(s, 457200, y, 1554600, 822900, fill, border, 12700)
    add_box_text(s, label, 457200, y + 100000, 1554600, 300000,
                 font_size=MEDIUM_SIZE, color=DARK_TEXT, bold=True)
    add_box_text(s, name, 2194560, y + 50000, 6492240, 250000,
                 font_size=BODY_SIZE, color=DARK_TEXT, bold=True,
                 alignment=PP_ALIGN.LEFT, font_name="Calibri")
    add_label(s, desc, 2194560, y + 300000, 6492240, 500000,
              font_size=SMALL_SIZE, color=GRAY_TEXT, italic=False,
              alignment=PP_ALIGN.LEFT, font_name="Calibri")
    y += 930000

# Hardware needs at bottom
add_body_text(s, 457200, y + 50000, 8229600, 400000, [
    ("Hardware needed: ", DARK_TEXT, True, False),
    ("Global shutter USB camera (~$30) | 3D-printed bracket | 170mm paddle disc | Orange ping-pong balls", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# SLIDE 12: Demo Videos (placeholder for live presentation)
# ============================================================
s = add_slide_default()
add_title(s, "Demo Videos")

add_label(s, "RSL-RL PPO  |  Isaac Lab + PhysX 5  |  12,288 parallel envs",
          457200, 731520, 8229600, 228600,
          font_size=BODY_SIZE, color=GRAY_TEXT, italic=False)

# Two placeholder cards for videos
add_box(s, 457200, 1200000, 3931800, 2500000, LIGHT_GRAY_BG, GRAY_BORDER, 12700)
add_label(s, "Pi2: Pose-only torso tracking\n(no velocity commands)",
          457200, 2200000, 3931800, 400000,
          font_size=BODY_SIZE, color=MED_GRAY, italic=True)
add_label(s, "[Insert video / screenshot]",
          457200, 1500000, 3931800, 400000,
          font_size=SUBTITLE_SIZE, color=GRAY_TEXT, italic=True)

add_box(s, 4754880, 1200000, 3931800, 2500000, LIGHT_GRAY_BG, GRAY_BORDER, 12700)
add_label(s, "Pi2: Full 8D tracking\n(walking + pose + velocity)",
          4754880, 2200000, 3931800, 400000,
          font_size=BODY_SIZE, color=MED_GRAY, italic=True)
add_label(s, "[Insert video / screenshot]",
          4754880, 1500000, 3931800, 400000,
          font_size=SUBTITLE_SIZE, color=GRAY_TEXT, italic=True)

# Captions
add_body_text(s, 457200, 3900000, 8229600, 800000, [
    ("Video locations:", DARK_TEXT, True, False),
    ("  videos/pi2_no_velocity/ — pose-only tracking (with and without exponential smoothing)", GRAY_TEXT),
    ("  videos/pi2_velocity_tracking/ — current training run with full 8D commands", GRAY_TEXT),
    ("  myrecordings/ — raw screen recordings (Feb 26 − Mar 11)", GRAY_TEXT),
], font_size=SMALL_SIZE)


# ============================================================
# Save
# ============================================================
out_path = "slides/weekmarch12thcapstone.pptx"
prs.save(out_path)
print(f"Saved to {out_path}")
print(f"Total slides: {len(prs.slides)}")
